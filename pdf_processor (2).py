from fastapi import APIRouter, HTTPException, UploadFile, File, Depends
from pydantic import BaseModel
import requests
import os
import json
import random
import io
import time
import math
from typing import Union, List, Dict
from bs4 import BeautifulSoup
import re
from docx import Document
from PIL import Image
import pytesseract
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from datetime import datetime
from urllib.parse import urlparse, parse_qs
from usage_limiter import enforce_usage_limit
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()

# --- Configuration Constants ---
TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY")
TOGETHER_API_URL = "https://api.together.xyz/v1/chat/completions"
TOGETHER_MODEL = "mistralai/Mixtral-8x7B-Instruct-v0.1"
MAX_LLM_INPUT_CHARS = 28000
YOUTUBE_API_KEY = os.getenv("YOUTUBE_API_KEY")
MAX_FILE_SIZE_MB = 10
MAX_RETRIES = 3
RETRY_DELAY = 1

# Validate required environment variables
if not TOGETHER_API_KEY:
    raise ValueError("TOGETHER_API_KEY environment variable not set")

HEADERS = {
    "Authorization": f"Bearer {TOGETHER_API_KEY}",
    "Content-Type": "application/json",
    "Accept": "application/json"
}

# --- Pydantic Models ---
class ContentRequest(BaseModel):
    text: str

class GenerateQuestionsRequest(BaseModel):
    text: str
    difficulty: str = "medium"
    count: int = 5

class FollowUpRequest(BaseModel):
    summary: str
    question: str

class URLRequest(BaseModel):
    url: str

class YouTubeURLRequest(BaseModel):
    youtube_url: str

class FlashcardsRequest(BaseModel):
    text: str

class VocabularyRequest(BaseModel):
    text: str

class HumanizeRequest(BaseModel):
    text: str

class MindMapRequest(BaseModel):
    text: str

class DiagramRequest(BaseModel):
    text: str
    diagram_type: str = "flowchart"  # flowchart, sequence, class, state, entity

class HandwrittenRequest(BaseModel):
    text: str
    style: str = "neat"  # neat, casual, messy

class FileUploadResponse(BaseModel):
    extracted_text: str
    file_type: str
    file_size: int

class QuestionItem(BaseModel):
    text: str
    options: List[str]
    answer: str

class FlashcardItem(BaseModel):
    front: str
    back: str

class VocabularyItem(BaseModel):
    word: str
    definition: str

# --- Helper Functions ---
def validate_file_size(file: UploadFile) -> None:
    max_size = MAX_FILE_SIZE_MB * 1024 * 1024
    file.file.seek(0, 2)  # Seek to end
    file_size = file.file.tell()
    file.file.seek(0)  # Reset pointer
    if file_size > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"File size exceeds maximum allowed {MAX_FILE_SIZE_MB}MB"
        )

async def call_together_ai(
    prompt_messages: List[Dict[str, str]],
    temperature: float = 0.7,
    max_tokens: int = 1024,
    retries: int = MAX_RETRIES
) -> str:
    payload = {
        "model": TOGETHER_MODEL,
        "messages": prompt_messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
        "top_p": 0.9,
        "top_k": 50,
        "repetition_penalty": 1.0,
        "stop": ["</s>"]
    }

    for attempt in range(retries):
        try:
            logger.info(f"Calling Together AI API (attempt {attempt + 1})")
            start_time = time.time()
            
            response = requests.post(
                TOGETHER_API_URL,
                headers=HEADERS,
                json=payload,
                timeout=30
            )
            response.raise_for_status()
            
            result = response.json()
            if not result or "choices" not in result or not result["choices"]:
                raise ValueError("Invalid response structure from Together AI")
                
            elapsed_time = time.time() - start_time
            logger.info(f"API call completed in {elapsed_time:.2f}s")
            
            return result["choices"][0]["message"]["content"]
        
        except requests.exceptions.HTTPError as e:
            error_msg = f"HTTP Error: {e.response.status_code} - {e.response.text}"
            logger.error(error_msg)
            
            if e.response.status_code == 429:
                wait_time = min((2 ** attempt) * RETRY_DELAY, 60)
                logger.warning(f"Rate limited. Waiting {wait_time}s before retry...")
                time.sleep(wait_time)
                continue
                
            raise HTTPException(
                status_code=e.response.status_code,
                detail=f"Together AI API error: {error_msg}"
            )
            
        except requests.exceptions.RequestException as e:
            logger.error(f"Request failed: {str(e)}")
            if attempt == retries - 1:
                raise HTTPException(
                    status_code=503,
                    detail="Service temporarily unavailable"
                )
            time.sleep(RETRY_DELAY)
            
        except Exception as e:
            logger.error(f"Unexpected error: {str(e)}")
            if attempt == retries - 1:
                raise HTTPException(
                    status_code=500,
                    detail=f"Failed to process AI request: {str(e)}"
                )
            time.sleep(RETRY_DELAY)

    raise HTTPException(
        status_code=500,
        detail="Failed to call Together AI after multiple attempts"
    )

async def extract_text_from_pdf(pdf_file: UploadFile) -> str:
    try:
        logger.info(f"Extracting text from PDF: {pdf_file.filename}")
        pdf_bytes = await pdf_file.read()
        
        if not pdf_bytes:
            raise ValueError("Empty PDF file")
            
        text = pdf_extract_text(io.BytesIO(pdf_bytes))
        if not text.strip():
            raise ValueError("No readable text found in PDF")
            
        return text
        
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract text from PDF: {str(e)}"
        )

async def extract_text_from_docx(docx_file: UploadFile) -> str:
    try:
        logger.info(f"Extracting text from DOCX: {docx_file.filename}")
        docx_bytes = await docx_file.read()
        document = Document(io.BytesIO(docx_bytes))
        
        full_text = []
        for para in document.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
                
        extracted_text = '\n'.join(full_text)
        
        if not extracted_text.strip():
            raise ValueError("No readable text found in DOCX")
            
        return extracted_text
        
    except Exception as e:
        logger.error(f"DOCX extraction error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract text from Word document: {str(e)}"
        )

async def extract_text_from_ppt(ppt_file: UploadFile) -> str:
    try:
        logger.info(f"Extracting text from PPT: {ppt_file.filename}")
        ppt_bytes = await ppt_file.read()
        prs = Presentation(io.BytesIO(ppt_bytes))
        
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)
                    
        extracted_text = '\n'.join(full_text)
        
        if not extracted_text.strip():
            raise ValueError("No readable text found in PPT")
            
        return extracted_text
        
    except Exception as e:
        logger.error(f"PPT extraction error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract text from PowerPoint: {str(e)}"
        )

async def extract_text_from_image(image_file: UploadFile) -> str:
    try:
        logger.info(f"Extracting text from image: {image_file.filename}")
        image_bytes = await image_file.read()
        image = Image.open(io.BytesIO(image_bytes))
        
        # Configure Tesseract (if needed)
        custom_config = r'--oem 3 --psm 6'
        text = pytesseract.image_to_string(image, config=custom_config)
        
        if not text.strip():
            raise ValueError("No text found in image via OCR")
            
        return text
        
    except Exception as e:
        logger.error(f"Image OCR error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract text from image: {str(e)}"
        )


def clean_extracted_text(text: str) -> str:
    """Clean and normalize extracted text"""
    if not text:
        return ""
        
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    
    # Remove non-printable characters except basic punctuation
    text = re.sub(r'[^\x20-\x7E\u2018\u2019\u201C\u201D\u2013\u2014]', ' ', text)
    
    # Truncate if too long
    if len(text) > MAX_LLM_INPUT_CHARS:
        text = text[:MAX_LLM_INPUT_CHARS]
        logger.warning(f"Text truncated to {MAX_LLM_INPUT_CHARS} characters")
        
    return text

# --- API Endpoints ---
@router.post("/api/upload-and-extract",
             response_model=FileUploadResponse,
             response_description="Extracted text from uploaded file")
async def upload_and_extract(
    file: UploadFile = File(...),
    user_id: str = enforce_usage_limit("uploads")
) -> FileUploadResponse:
    """
    Handles file uploads (PDF, DOCX, PPTX, Image) and extracts text content.
    Automatically detects file type and calls the appropriate extractor.
    """
    try:
        validate_file_size(file)
        filename = file.filename.lower()
        file_type = filename.split('.')[-1]
        
        logger.info(f"Processing file upload: {filename} for user {user_id}")
        
        if filename.endswith('.pdf'):
            text = await extract_text_from_pdf(file)
        elif filename.endswith('.docx'):
            text = await extract_text_from_docx(file)
        elif filename.endswith(('.ppt', '.pptx')):
            text = await extract_text_from_ppt(file)
        elif filename.endswith(('.jpg', '.jpeg', '.png')):
            text = await extract_text_from_image(file)
        else:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Supported: PDF, DOCX, PPT/PPTX, JPG/PNG"
            )
            
        cleaned_text = clean_extracted_text(text)
        
        # Get file size in bytes
        await file.seek(0, 2)
        file_size = file.tell()
        await file.seek(0)
        
        return {
            "extracted_text": cleaned_text,
            "file_type": file_type,
            "file_size": file_size
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File upload processing error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process uploaded file: {str(e)}"
        )

@router.post("/api/fetch-and-extract-url",
             response_model=ContentRequest,
             response_description="Extracted text from web page")
async def fetch_and_extract_url(
    request: URLRequest,
    user_id: str = enforce_usage_limit("summaries")
) -> ContentRequest:
    """
    Fetches content from a given web page URL and extracts its main readable text.
    Uses BeautifulSoup to parse HTML and remove irrelevant elements.
    """
    url = request.url.strip()
    if not url:
        raise HTTPException(status_code=400, detail="No URL provided")
        
    # Prepend 'https://' if no scheme is provided
    if not url.startswith(('http://', 'https://')):
        url = 'https://' + url

    try:
        logger.info(f"Fetching URL content: {url} for user {user_id}")
        
        response = requests.get(url, timeout=10, headers={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })
        response.raise_for_status()

        soup = BeautifulSoup(response.text, 'lxml')
        
        # Remove unwanted elements
        for element in soup(["script", "style", "nav", "footer", "header", "form", 
                           "aside", "meta", "iframe", "img", "svg", "link", 
                           "input", "button", "select", "textarea"]):
            element.decompose()

        # Try to find main content areas
        content_selectors = [
            'article', 'main', 'div.main-content', 'div.entry-content', 
            'div.post-content', 'div.article-body', 'div[role="main"]',
            'div#content', 'div#main', 'div.content', 'div.post',
            'div.blog-post', 'div.article'
        ]
        
        main_text = ""
        for selector in content_selectors:
            elements = soup.select(selector)
            for element in elements:
                text = element.get_text(separator=' ', strip=True)
                if len(text) > len(main_text):
                    main_text = text

        # Fallback to body text if no main content found
        if not main_text or len(main_text) < 100:
            body = soup.find('body')
            if body:
                main_text = body.get_text(separator=' ', strip=True)

        # Final fallback to entire document
        if not main_text or len(main_text) < 50:
            main_text = soup.get_text(separator=' ', strip=True)

        cleaned_text = clean_extracted_text(main_text)
        
        if not cleaned_text:
            raise HTTPException(
                status_code=400,
                detail="Could not extract meaningful text from the URL"
            )
            
        return {"text": cleaned_text}
        
    except requests.exceptions.RequestException as e:
        logger.error(f"URL fetch error: {str(e)}")
        if isinstance(e, requests.exceptions.Timeout):
            raise HTTPException(status_code=408, detail="Request to URL timed out")
        elif isinstance(e, requests.exceptions.ConnectionError):
            raise HTTPException(status_code=503, detail="Could not connect to URL")
        elif hasattr(e, 'response') and e.response:
            raise HTTPException(
                status_code=e.response.status_code,
                detail=f"Error fetching URL: {e.response.status_code} {e.response.reason}"
            )
        raise HTTPException(status_code=500, detail=f"Failed to fetch URL: {str(e)}")
        
    except Exception as e:
        logger.error(f"URL processing error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process URL content: {str(e)}"
        )


@router.post("/api/summarize",
             response_model=ContentRequest,
             response_description="Generated summary of the input text")
async def summarize_text(
    request: ContentRequest,
    user_id: str = enforce_usage_limit("summaries")
) -> ContentRequest:
    """
    Generates a detailed, structured summary from the provided text using an LLM.
    The summary includes key points and maintains the original meaning.
    """
    try:
        logger.info(f"Generating summary for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """You are an expert at creating detailed, structured summaries. 
                Use markdown formatting with headings (##) and bullet points.
                Include all key concepts and maintain the original meaning."""
            },
            {
                "role": "user",
                "content": f"""Create a comprehensive summary of this content:
                
                {request.text}"""
            }
        ]

        summary = await call_together_ai(prompt_messages, max_tokens=1024)
        return {"text": summary}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Summary generation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate summary: {str(e)}"
        )

@router.post("/api/generate-questions",
             response_model=List[QuestionItem],
             response_description="List of generated questions with options")
async def generate_questions(
    request: GenerateQuestionsRequest,
    user_id: str = enforce_usage_limit("questions")
) -> List[QuestionItem]:
    """
    Generates multiple-choice questions from the provided text.
    Each question includes 4 options and 1 correct answer.
    """
    try:
        logger.info(f"Generating {request.count} {request.difficulty} questions for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Generate multiple-choice questions in JSON format.
                Each question must have:
                - 'text': The question text
                - 'options': List of 4 options
                - 'answer': The correct answer"""
            },
            {
                "role": "user",
                "content": f"""Create {request.count} {request.difficulty} difficulty questions from this text.
                Format as a JSON list of question objects.
                
                Text:
                {request.text}"""
            }
        ]

        response = await call_together_ai(prompt_messages, max_tokens=1500)
        
        # Clean and validate response
        if response.startswith("```json"):
            response = response[7:-3].strip()
        
        questions = json.loads(response)
        if not isinstance(questions, list):
            raise ValueError("Invalid question format")
        
        # Validate and randomize each question
        validated_questions = []
        for q in questions:
            if not all(k in q for k in ["text", "options", "answer"]):
                logger.warning(f"Skipping invalid question: {q}")
                continue
                
            if len(q["options"]) != 4:
                logger.warning(f"Question has incorrect options count: {q['text']}")
                continue
                
            if q["answer"] not in q["options"]:
                logger.warning(f"Correct answer not in options for: {q['text']}")
                continue
                
            random.shuffle(q["options"])
            validated_questions.append(q)
        
        if not validated_questions:
            raise HTTPException(
                status_code=500,
                detail="No valid questions could be generated"
            )
            
        return validated_questions
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to parse generated questions: {str(e)}"
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Question generation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate questions: {str(e)}"
        )

@router.post("/api/follow-up",
             response_model=ContentRequest,
             response_description="Answer to the follow-up question")
async def follow_up(
    request: FollowUpRequest,
    user_id: str = enforce_usage_limit("summaries")
) -> ContentRequest:
    """
    Answers a follow-up question based on the provided summary.
    The response is generated using the context from the summary only.
    """
    try:
        logger.info(f"Processing follow-up question for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Answer questions concisely based ONLY on the provided summary.
                If the answer isn't in the summary, say 'I cannot answer based on the provided information'."""
            },
            {
                "role": "user",
                "content": f"""Summary:
                {request.summary}
                
                Question: {request.question}"""
            }
        ]

        answer = await call_together_ai(prompt_messages, max_tokens=300)
        return {"text": answer}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Follow-up processing failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to process follow-up question: {str(e)}"
        )

@router.post("/api/generate-flashcards",
             response_model=List[FlashcardItem],
             response_description="List of generated flashcards")
async def generate_flashcards(
    request: FlashcardsRequest,
    user_id: str = enforce_usage_limit("flashcards")
) -> List[FlashcardItem]:
    """
    Generates flashcards (front and back) from the provided text.
    Each flashcard contains a concept/question on the front and explanation/answer on the back.
    """
    try:
        logger.info(f"Generating flashcards for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Generate flashcards in JSON format with:
                - 'front': Concept or question
                - 'back': Definition or answer
                Return ONLY the JSON array of flashcards."""
            },
            {
                "role": "user",
                "content": f"""Create flashcards from this content:
                
                {request.text}"""
            }
        ]

        response = await call_together_ai(prompt_messages, max_tokens=1500)
        
        # Clean and validate response
        if response.startswith("```json"):
            response = response[7:-3].strip()
        
        flashcards = json.loads(response)
        if not isinstance(flashcards, list):
            raise ValueError("Invalid flashcard format")
        
        # Validate each flashcard
        validated_flashcards = []
        for card in flashcards:
            if not all(k in card for k in ["front", "back"]):
                logger.warning(f"Skipping invalid flashcard: {card}")
                continue
                
            if not card["front"].strip() or not card["back"].strip():
                logger.warning(f"Skipping empty flashcard: {card}")
                continue
                
            validated_flashcards.append(card)
        
        if not validated_flashcards:
            raise HTTPException(
                status_code=500,
                detail="No valid flashcards could be generated"
            )
            
        return validated_flashcards
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to parse generated flashcards: {str(e)}"
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Flashcard generation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate flashcards: {str(e)}"
        )

@router.post("/api/generate-vocabulary",
             response_model=List[VocabularyItem],
             response_description="List of vocabulary words with definitions")
async def generate_vocabulary(
    request: VocabularyRequest,
    user_id: str = enforce_usage_limit("vocabulary")
) -> List[VocabularyItem]:
    """
    Extracts important vocabulary words from the text along with their definitions.
    The definitions are contextually relevant to how the words are used in the text.
    """
    try:
        logger.info(f"Generating vocabulary list for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Extract vocabulary words with definitions in JSON format.
                Each item should have:
                - 'word': The vocabulary word
                - 'definition': Contextual definition
                Return ONLY the JSON array of vocabulary items."""
            },
            {
                "role": "user",
                "content": f"""Identify key vocabulary from this text:
                
                {request.text}"""
            }
        ]

        response = await call_together_ai(prompt_messages, max_tokens=1024)
        
        # Clean and validate response
        if response.startswith("```json"):
            response = response[7:-3].strip()
        
        vocabulary = json.loads(response)
        if not isinstance(vocabulary, list):
            raise ValueError("Invalid vocabulary format")
        
        # Validate each vocabulary item
        validated_vocab = []
        for item in vocabulary:
            if not all(k in item for k in ["word", "definition"]):
                logger.warning(f"Skipping invalid vocabulary item: {item}")
                continue
                
            if not item["word"].strip() or not item["definition"].strip():
                logger.warning(f"Skipping empty vocabulary item: {item}")
                continue
                
            validated_vocab.append(item)
        
        if not validated_vocab:
            raise HTTPException(
                status_code=500,
                detail="No vocabulary could be extracted"
            )
            
        return validated_vocab
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to parse vocabulary: {str(e)}"
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Vocabulary extraction failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to extract vocabulary: {str(e)}"
        )

@router.post("/api/humanize-text",
             response_model=ContentRequest,
             response_description="Humanized version of the input text")
async def humanize_text(
    request: HumanizeRequest,
    user_id: str = enforce_usage_limit("humanize")
) -> ContentRequest:
    """
    Rewrites the provided text to sound more natural and human-like,
    while maintaining the original meaning and key information.
    """
    try:
        logger.info(f"Humanizing text for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Rewrite text to sound natural and human-like.
                Use conversational tone, vary sentence structure, and add natural flow.
                Maintain all key information from the original."""
            },
            {
                "role": "user",
                "content": f"""Humanize this text:
                
                {request.text}"""
            }
        ]

        humanized = await call_together_ai(prompt_messages, temperature=0.8, max_tokens=1500)
        return {"text": humanized}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Text humanization failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to humanize text: {str(e)}"
        )

@router.post("/api/generate-mindmap",
             response_model=ContentRequest,
             response_description="Mermaid.js code for the generated mindmap")
async def generate_mindmap(
    request: MindMapRequest,
    user_id: str = enforce_usage_limit("diagrams")
) -> ContentRequest:
    """
    Generates a mind map structure from the provided text in Mermaid.js format.
    The mind map has a hierarchical structure with a clear root node.
    """
    try:
        logger.info(f"Generating mindmap for user {user_id}")
        
        prompt_messages = [
            {
                "role": "system",
                "content": """Generate mindmaps in Mermaid.js format.
                Start with 'mindmap' and use proper indentation.
                Include a root node and hierarchical structure.
                Return ONLY the Mermaid code, DO NOT INCLUDE ANY SPECIAL CHARACTERS."""
            },
            {
                "role": "user",
                "content": f"""Create a mindmap from this content:
                
                {request.text}"""
            }
        ]

        code = await call_together_ai(prompt_messages, temperature=0.3, max_tokens=1024)
        
        # Ensure valid Mermaid syntax
        if not code.strip().startswith("mindmap"):
            code = "mindmap\n  root((Main Topic))\n" + code
            
        return {"text": code}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Mindmap generation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate mindmap: {str(e)}"
        )

@router.post("/api/generate-diagram",
             response_model=ContentRequest,
             response_description="Mermaid.js code for the generated diagram")
async def generate_diagram(
    request: DiagramRequest,
    user_id: str = enforce_usage_limit("diagrams")
) -> ContentRequest:
    """
    Generates a diagram from the provided text in Mermaid.js format.
    The type of diagram (flowchart, sequence, etc.) is specified in the request.
    """
    try:
        logger.info(f"Generating {request.diagram_type} diagram for user {user_id}")
        
        diagram_types = {
            "flowchart": "flowchart TD",
            "sequence": "sequenceDiagram",
            "class": "classDiagram",
            "state": "stateDiagram-v2",
            "entity": "erDiagram"
        }
        
        if request.diagram_type not in diagram_types:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid diagram type. Supported: {', '.join(diagram_types.keys())}"
            )
            
        prompt_messages = [
            {
                "role": "system",
                "content": f"""Generate a {request.diagram_type} diagram in Mermaid.js format.
                Start with '{diagram_types[request.diagram_type]}'.
                Use proper syntax for the diagram type.
                Return ONLY the Mermaid code."""
            },
            {
                "role": "user",
                "content": f"""Create a {request.diagram_type} diagram from this content:
                
                {request.text}"""
            }
        ]

        code = await call_together_ai(prompt_messages, temperature=0.3, max_tokens=1024)
        
        # Ensure valid Mermaid syntax
        if not code.strip().startswith(diagram_types[request.diagram_type]):
            code = f"{diagram_types[request.diagram_type]}\n" + code
            
        return {"text": code}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Diagram generation failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate diagram: {str(e)}"
        )

@router.post("/api/generate-handwritten",
             response_model=ContentRequest,
             response_description="Text formatted to resemble handwritten notes")
async def generate_handwritten(
    request: HandwrittenRequest,
    user_id: str = enforce_usage_limit("handwritten")
) -> ContentRequest:
    """
    Converts the provided text into a format that resembles handwritten notes,
    with stylistic elements based on the requested style (neat, casual, messy).
    """
    try:
        logger.info(f"Generating {request.style} handwritten notes for user {user_id}")
        
        style_descriptions = {
            "neat": "neat and organized handwriting like careful notes",
            "casual": "casual everyday handwriting with natural variation",
            "messy": "quick, messy handwriting like lecture notes"
        }
        
        if request.style not in style_descriptions:
            raise HTTPException(
                status_code=400,
                detail="Invalid style. Supported: neat, casual, messy"
            )
            
        prompt_messages = [
            {
                "role": "system",
                "content": f"""Convert text to {style_descriptions[request.style]}.
                Use markdown to represent handwritten features:
                - ~crossed out~ text
                - **underlined** terms
                - [margin notes in brackets]
                - [doodle: description] for illustrations"""
            },
            {
                "role": "user",
                "content": f"""Convert this to {request.style} handwriting:
                
                {request.text}"""
            }
        ]

        notes = await call_together_ai(prompt_messages, temperature=0.7, max_tokens=1500)
        return {"text": notes}
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Handwritten conversion failed: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate handwritten notes: {str(e)}"
        )